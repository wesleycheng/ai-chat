#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
AI Chat Platform - 专业版项目文档 PPT 生成器
包含图文混排、专业配色、架构图等元素
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt
import pptx


def create_professional_presentation():
    """创建专业版项目演示 PPT"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # 定义专业配色方案
    COLOR_PRIMARY = RGBColor(0x1F, 0x4F, 0x8F)      # 深蓝色
    COLOR_SECONDARY = RGBColor(0x3B, 0x82, 0xF6)    # 亮蓝色
    COLOR_ACCENT = RGBColor(0xF5, 0x93, 0x1E)       # 橙色
    COLOR_SUCCESS = RGBColor(0x10, 0xB9, 0x81)      # 绿色
    COLOR_DARK = RGBColor(0x1F, 0x29, 0x3B)        # 深灰色
    COLOR_LIGHT_GRAY = RGBColor(0xF1, 0xF5, 0xF9)  # 浅灰色
    COLOR_TEXT = RGBColor(0x33, 0x33, 0x33)         # 正文灰色
    
    # 幻灯片 1: 专业封面
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    add_professional_title_slide(
        slide1, 
        "AI Chat Platform",
        "Enterprise Multi-Modal AI Chat Platform\n项目开发与技术文档",
        COLOR_PRIMARY,
        COLOR_SECONDARY
    )
    add_text_box(slide1, "Prepared by: Wesley Cheng", 7.5, 6.8, 2.0, 0.4, 
                  size=14, color=COLOR_DARK, align=PP_ALIGN.RIGHT)
    add_text_box(slide1, "Date: May 30, 2026", 7.5, 7.1, 2.0, 0.4,
                  size=14, color=COLOR_DARK, align=PP_ALIGN.RIGHT)
    
    # 幻灯片 2: 项目概述（图文版）
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide2, "项目概述", COLOR_PRIMARY, COLOR_LIGHT_GRAY)
    
    # 左侧：项目信息表格
    add_info_card(slide2, 0.5, 1.2, 4.5, 2.8, 
                  "📊 项目信息", 
                  ["项目名称: AI Chat Platform",
                   "开发周期: 2026-05-27 ~ 2026-05-30",
                   "总提交次数: 28 commits",
                   "代码规模: 前端~15K行 / 后端~8K行",
                   "生产环境: http://8.137.103.202"],
                  COLOR_SECONDARY)
    
    # 右侧：核心特性
    add_info_card(slide2, 5.2, 1.2, 4.3, 2.8,
                  "✨ 核心特性",
                  ["✓ 多模型支持 (DeepSeek/OpenAI/Anthropic/Ollama)",
                   "✓ 智能文件解析 (PDF/DOCX/XLSX/TXT/MD)",
                   "✓ Agent管理系统 (自定义Prompt与工具)",
                   "✓ 流式对话 (SSE实时输出)",
                   "✓ 响应式设计 (桌面端+移动端)",
                   "✓ 安全认证 (JWT + bcrypt)"],
                  COLOR_SUCCESS)
    
    # 底部：技术亮点
    add_highlight_bar(slide2, 0.5, 4.5, 9.0, 0.8, 
                      "🚀 技术亮点: 前后端分离 | Docker容器化 | RAG检索增强 | 异步任务队列",
                      COLOR_ACCENT)
    
    # 幻灯片 3: 系统架构（视觉化）
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide3, "系统架构设计", COLOR_PRIMARY, COLOR_LIGHT_GRAY)
    
    # 绘制架构层次图
    add_architecture_diagram(slide3, COLOR_PRIMARY, COLOR_SECONDARY, COLOR_SUCCESS)
    
    # 幻灯片 4: 技术栈（分栏展示）
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide4, "技术栈", COLOR_PRIMARY, COLOR_LIGHT_GRAY)
    
    # 前端技术栈
    add_tech_stack_card(slide4, 0.5, 1.2, 4.5, 5.0, 
                        "🎨 前端技术栈",
                        ["React 18.3 + TypeScript 5.5",
                         "Vite 5.4 (构建工具)",
                         "Tailwind CSS 3.4 (样式)",
                         "Zustand 4.5 (状态管理)",
                         "@tanstack/react-query 5.56",
                         "React Router DOM 6.26",
                         "Axios 1.7 (HTTP客户端)",
                         "react-markdown 9.0",
                         "lucide-react 0.400 (图标)"],
                        COLOR_SECONDARY)
    
    # 后端技术栈
    add_tech_stack_card(slide4, 5.2, 1.2, 4.3, 5.0,
                        "⚙️ 后端技术栈",
                        ["FastAPI 0.115 + Uvicorn",
                         "SQLAlchemy 2.0 + Alembic",
                         "asyncpg 0.30 + psycopg2-binary",
                         "python-jose 3.3 + passlib 1.7",
                         "LangChain 0.3 + LangGraph 0.2",
                         "ChromaDB 0.5 (向量数据库)",
                         "arq 0.24 + Redis 5.2",
                         "PyMuPDF 1.25 + python-docx 1.1"],
                        COLOR_SUCCESS)
    
    # 幻灯片 5: 功能模块（图标化展示）
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide5, "功能模块", COLOR_PRIMARY, COLOR_LIGHT_GRAY)
    
    # 5个功能模块的卡片展示
    modules = [
        ("1", "用户认证模块", ["JWT Token认证", "用户注册/登录", "密码bcrypt加密", "Token刷新机制"]),
        ("2", "模型配置模块", ["多模型支持", "CRUD管理", "API Key加密", "默认模型设置"]),
        ("3", "对话管理模块", ["SSE流式输出", "上下文管理", "Agent选择", "多轮对话"]),
        ("4", "文件管理模块", ["多格式支持", "异步上传", "自动解析", "RAG集成"]),
        ("5", "Agent管理模块", ["自定义Prompt", "工具配置", "CRUD管理", "模型关联"])
    ]
    
    x_positions = [0.5, 2.3, 4.1, 5.9, 7.7]
    for i, (num, title, features) in enumerate(modules):
        add_module_card(slide5, x_positions[i], 1.5, 1.6, 4.5, 
                        num, title, features, COLOR_SECONDARY if i % 2 == 0 else COLOR_SUCCESS)
    
    # 幻灯片 6: 开发过程（时间线）
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide6, "开发过程", COLOR_PRIMARY, COLOR_LIGHT_GRAY)
    
    # 时间线展示
    add_timeline(slide6, 
                 [(0.8, "Day 1\n05-27", ["项目初始化", "Git仓库搭建", "Docker环境配置", "前后端骨架"]),
                  (2.8, "Day 2\n05-28", ["用户认证系统", "模型配置CRUD", "SSE流式输出", "ECS部署"]),
                  (4.8, "Day 3\n05-29", ["Agent管理页面", "文件上传功能", "Agent对话集成", "响应式优化"]),
                  (6.8, "Day 4\n05-30", ["默认模型编辑", "输入优化", "加载动画", "文档生成"])],
                 COLOR_SECONDARY)
    
    # 幻灯片 7: 提交日志（表格化）
    slide7 = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide7, "提交日志统计", COLOR_PRIMARY, COLOR_LIGHT_GRAY)
    
    # 提交统计
    add_commit_statistics(slide7, COLOR_SUCCESS)
    
    # 最近10条提交记录（表格）
    commits = [
        ("907c52b", "2026-05-30", "feat", "设置页面支持编辑默认模型和启用状态"),
        ("b23c987", "2026-05-29", "fix", "修复文件上传问题，修复认证头拼写错误"),
        ("eedd3dd", "2026-05-29", "feat", "优化所有页面响应式设计"),
        ("adcb7fb", "2026-05-29", "feat", "响应式UI优化"),
        ("c6a6c10", "2026-05-29", "fix", "修复文件上传问题"),
        ("19dba34", "2026-05-29", "fix", "修复文件内容获取问题"),
        ("fbf714c", "2026-05-29", "fix", "优化文件上传流程"),
        ("29ec3ab", "2026-05-29", "feat", "聊天支持上传文件"),
        ("0b76cfd", "2026-05-29", "feat", "聊天页面支持选择Agent"),
        ("9e72e15", "2026-05-29", "fix", "修复创建Agent时空字符串错误")
    ]
    add_commit_table(slide7, commits, COLOR_SECONDARY)
    
    # 幻灯片 8: Bug修复记录（问题-解决方案对照）
    slide8 = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide8, "重大Bug修复记录", COLOR_PRIMARY, COLOR_LIGHT_GRAY)
    
    bugs = [
        ("Bug #1", "Authorization拼写错误", 
         "前端api.ts中Authorization拼写为Authoriation",
         "修正拼写错误", "b23c987", COLOR_ACCENT),
        ("Bug #2", "文件上传失败",
         "手动设置Content-Type导致boundary缺失",
         "删除手动Content-Type设置", "c6a6c10", COLOR_ACCENT),
        ("Bug #3", "Agent创建500错误",
         "空字符串''导致数据库外键约束报错",
         "将空字符串转为None", "9e72e15", COLOR_ACCENT),
        ("Bug #4", "文件内容未随消息发送",
         "file_ids类型不匹配，查询条件错误",
         "转换file_ids为UUID列表", "19dba34", COLOR_ACCENT)
    ]
    
    for i, (bug_id, title, problem, solution, commit, color) in enumerate(bugs):
        add_bug_fix_card(slide8, 0.5, 1.3 + i*1.3, 9.0, 1.1,
                         bug_id, title, problem, solution, commit, color)
    
    # 幻灯片 9: 部署架构（可视化）
    slide9 = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide9, "部署架构", COLOR_PRIMARY, COLOR_LIGHT_GRAY)
    
    # Docker服务架构图
    add_deployment_architecture(slide9, COLOR_SECONDARY, COLOR_SUCCESS)
    
    # 环境变量配置
    add_info_card(slide9, 0.5, 5.5, 9.0, 1.5,
                  "🔧 环境变量配置",
                  ["DATABASE_URL - 数据库连接字符串",
                   "REDIS_URL - Redis连接字符串",
                   "CHROMA_HOST - ChromaDB主机地址",
                   "SECRET_KEY - JWT密钥"],
                  COLOR_ACCENT)
    
    # 幻灯片 10: 使用指南 & 未来规划
    slide10 = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide10, "使用指南 & 未来规划", COLOR_PRIMARY, COLOR_LIGHT_GRAY)
    
    # 使用指南
    add_info_card(slide10, 0.5, 1.2, 4.5, 2.8,
                  "📖 快速开始",
                  ["1. 访问 http://8.137.103.202",
                   "2. 注册账号并登录",
                   "3. 配置模型（设置->添加模型）",
                   "4. 开始对话（选择模型/Agent）",
                   "5. 上传文件进行RAG对话"],
                  COLOR_SECONDARY)
    
    # 未来规划
    add_info_card(slide10, 5.2, 1.2, 4.3, 2.8,
                  "🚀 未来规划",
                  ["功能增强:",
                   "  • 支持更多文件类型(PPT/MP4)",
                   "  • 实现RAG检索增强",
                   "  • 多模态输入(图片/语音)",
                   "性能优化:",
                   "  • 消息分页加载",
                   "  • 文件解析速度优化"],
                  COLOR_SUCCESS)
    
    # 底部：感谢页
    add_highlight_bar(slide10, 0.5, 4.5, 9.0, 0.8,
                      "🎯 项目已部署生产环境，欢迎测试与使用！",
                      COLOR_ACCENT)
    
    # 保存文件
    output_path = "/Users/wesleycheng/Documents/ai-chat-platform/AI_Chat_Platform_专业版项目文档.pptx"
    prs.save(output_path)
    print("✅ 专业版PPT生成成功！文件路径：{}".format(output_path))
    return output_path


# ============ 辅助函数 ============

def add_professional_title_slide(slide, title, subtitle, color_primary, color_secondary):
    """添加专业封面"""
    # 背景矩形
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(7.5)
    )
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0xF8, 0xFA, 0xFC)
    
    # 顶部装饰条
    top_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(0.1)
    )
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = color_primary
    top_bar.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.0), Inches(8), Inches(1.0))
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = color_primary
    p.alignment = PP_ALIGN.CENTER
    
    # 副标题
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.2), Inches(8), Inches(1.2))
    tf2 = subtitle_box.text_frame
    tf2.text = subtitle
    p2 = tf2.paragraphs[0]
    p2.font.size = Pt(18)
    p2.font.color.rgb = RGBColor(0x64, 0x74, 0x8B)
    p2.alignment = PP_ALIGN.CENTER


def add_slide_header(slide, title, color_primary, color_bg):
    """添加幻灯片标题栏"""
    # 背景条
    header_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(0.8)
    )
    header_bar.fill.solid()
    header_bar.fill.fore_color.rgb = color_primary
    header_bar.line.fill.background()
    
    # 标题文字
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.LEFT


def add_info_card(slide, x, y, width, height, title, items, color):
    """添加信息卡片"""
    # 卡片背景
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y),
        Inches(width), Inches(height)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    card.line.color.rgb = color
    card.line.width = Pt(1.5)
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.15), 
                                          Inches(width - 0.4), Inches(0.4))
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = color
    
    # 内容
    content_y = y + 0.6
    for item in items:
        item_box = slide.shapes.add_textbox(Inches(x + 0.3), Inches(content_y), 
                                             Inches(width - 0.6), Inches(0.35))
        tf2 = item_box.text_frame
        tf2.text = item
        p2 = tf2.paragraphs[0]
        p2.font.size = Pt(10)
        p2.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        content_y += 0.35


def add_text_box(slide, text, x, y, width, height, size=12, 
                 color=RGBColor(0, 0, 0), align=PP_ALIGN.LEFT):
    """添加文本框"""
    text_box = slide.shapes.add_textbox(Inches(x), Inches(y), 
                                         Inches(width), Inches(height))
    tf = text_box.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.alignment = align


def add_highlight_bar(slide, x, y, width, height, text, color):
    """添加高亮条"""
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y),
        Inches(width), Inches(height)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    
    text_box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.1), 
                                         Inches(width - 0.4), Inches(height - 0.2))
    tf = text_box.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.CENTER


def add_architecture_diagram(slide, color1, color2, color3):
    """添加架构图"""
    # 用户层
    add_arch_box(slide, 3.5, 1.2, 3.0, 0.6, "用户层 (Web Browser)", color1)
    
    # 箭头
    add_arrow(slide, 5.0, 1.8, 5.0, 2.2, color1)
    
    # 前端层
    add_arch_box(slide, 2.5, 2.2, 5.0, 0.6, "前端层 (React + TypeScript)", color2)
    
    # 箭头
    add_arrow(slide, 5.0, 2.8, 5.0, 3.2, color2)
    
    # 后端层
    add_arch_box(slide, 2.5, 3.2, 5.0, 0.6, "后端层 (FastAPI + Uvicorn)", color2)
    
    # 箭头
    add_arrow(slide, 5.0, 3.8, 5.0, 4.2, color2)
    
    # 数据层
    add_arch_box(slide, 2.5, 4.2, 5.0, 0.6, "数据层 (PostgreSQL + Redis + ChromaDB)", color3)
    
    # 技术特点（右侧）
    features = ["前后端分离", "Docker容器化", "异步任务(arq)", "向量检索(ChromaDB)", "流式输出(SSE)"]
    for i, feature in enumerate(features):
        feature_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(7.8), Inches(1.5 + i * 0.7),
            Inches(1.8), Inches(0.5)
        )
        feature_box.fill.solid()
        feature_box.fill.fore_color.rgb = color3
        feature_box.line.fill.background()
        
        tf = feature_box.text_frame
        tf.text = feature
        p = tf.paragraphs[0]
        p.font.size = Pt(9)
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.alignment = PP_ALIGN.CENTER


def add_arch_box(slide, x, y, width, height, text, color):
    """添加架构框"""
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y),
        Inches(width), Inches(height)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.line.color.rgb = color
    box.line.width = Pt(2)
    
    tf = box.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.CENTER


def add_connector(slide, x1, y1, x2, y2, color):
    """添加连接线"""
    connector = slide.shapes.add_connector(
        1,  # 直线
        Inches(x1), Inches(y1),
        Inches(x2), Inches(y2)
    )
    connector.line.color.rgb = color
    connector.line.width = Pt(2)


def add_tech_stack_card(slide, x, y, width, height, title, tech_list, color):
    """添加技术栈卡片"""
    # 卡片背景
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y),
        Inches(width), Inches(height)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    card.line.color.rgb = color
    card.line.width = Pt(2)
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.15), 
                                          Inches(width - 0.4), Inches(0.4))
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = color
    
    # 技术列表
    content_y = y + 0.6
    for tech in tech_list:
        tech_box = slide.shapes.add_textbox(Inches(x + 0.25), Inches(content_y), 
                                             Inches(width - 0.5), Inches(0.32))
        tf2 = tech_box.text_frame
        tf2.text = "• " + tech
        p2 = tf2.paragraphs[0]
        p2.font.size = Pt(9)
        p2.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        content_y += 0.32


def add_module_card(slide, x, y, width, height, num, title, features, color):
    """添加功能模块卡片"""
    # 卡片背景
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y),
        Inches(width), Inches(height)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    card.line.color.rgb = color
    card.line.width = Pt(1.5)
    
    # 模块编号（圆形）
    num_circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(x + width/2 - 0.25), Inches(y + 0.15),
        Inches(0.5), Inches(0.5)
    )
    num_circle.fill.solid()
    num_circle.fill.fore_color.rgb = color
    num_circle.line.fill.background()
    
    tf_num = num_circle.text_frame
    tf_num.text = num
    p_num = tf_num.paragraphs[0]
    p_num.font.size = Pt(16)
    p_num.font.bold = True
    p_num.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p_num.alignment = PP_ALIGN.CENTER
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(x + 0.1), Inches(y + 0.75), 
                                          Inches(width - 0.2), Inches(0.35))
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.CENTER
    
    # 功能列表
    content_y = y + 1.15
    for feature in features:
        feat_box = slide.shapes.add_textbox(Inches(x + 0.15), Inches(content_y), 
                                             Inches(width - 0.3), Inches(0.28))
        tf2 = feat_box.text_frame
        tf2.text = "✓ " + feature
        p2 = tf2.paragraphs[0]
        p2.font.size = Pt(8)
        p2.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        content_y += 0.28


def add_timeline(slide, milestones, color):
    """添加时间线"""
    for i, (x, label, items) in enumerate(milestones):
        # 时间点圆圈
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(x), Inches(1.5),
            Inches(0.6), Inches(0.6)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.line.fill.background()
        
        # 标签
        label_box = slide.shapes.add_textbox(Inches(x - 0.2), Inches(1.5), 
                                              Inches(1.0), Inches(0.6))
        tf = label_box.text_frame
        tf.text = label
        p = tf.paragraphs[0]
        p.font.size = Pt(9)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.alignment = PP_ALIGN.CENTER
        
        # 连接线
        if i < len(milestones) - 1:
            connector = slide.shapes.add_connector(
                1,  # 直线连接器
                Inches(x + 0.6), Inches(1.8),
                Inches(milestones[i+1][0]), Inches(1.8)
            )
            connector.line.color.rgb = color
            connector.line.width = Pt(2)
        
        # 内容卡片
        content_y = 2.3
        for item in items:
            item_box = slide.shapes.add_textbox(Inches(x - 0.3), Inches(content_y), 
                                                 Inches(1.8), Inches(0.35))
            tf2 = item_box.text_frame
            tf2.text = "✅ " + item
            p2 = tf2.paragraphs[0]
            p2.font.size = Pt(8)
            p2.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
            content_y += 0.35


def add_commit_statistics(slide, color):
    """添加提交统计"""
    stats = [
        ("✨ feat", "15次", color),
        ("🐛 fix", "12次", RGBColor(0xEF, 0x44, 0x44)),
        ("📝 docs", "1次", RGBColor(0x8B, 0x5C, 0xF6)),
        ("🔧 chore", "0次", RGBColor(0xF5, 0x93, 0x1E))
    ]
    
    for i, (label, count, color_stat) in enumerate(stats):
        stat_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5 + i * 2.3), Inches(1.2),
            Inches(2.0), Inches(0.6)
        )
        stat_box.fill.solid()
        stat_box.fill.fore_color.rgb = color_stat
        stat_box.line.fill.background()
        
        tf = stat_box.text_frame
        tf.text = label + ": " + count
        p = tf.paragraphs[0]
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.alignment = PP_ALIGN.CENTER


def add_commit_table(slide, commits, color):
    """添加提交记录表格"""
    # 表头
    headers = ["Commit Hash", "日期", "类型", "描述"]
    col_widths = [1.5, 1.2, 0.8, 4.5]
    col_x = [0.5, 2.0, 2.8, 3.6]
    
    # 表头背景
    header_bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(2.0),
        Inches(9.0), Inches(0.4)
    )
    header_bg.fill.solid()
    header_bg.fill.fore_color.rgb = color
    header_bg.line.fill.background()
    
    # 表头文字
    for i, header in enumerate(headers):
        header_box = slide.shapes.add_textbox(Inches(col_x[i]), Inches(2.0), 
                                              Inches(col_widths[i]), Inches(0.4))
        tf = header_box.text_frame
        tf.text = header
        p = tf.paragraphs[0]
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.alignment = PP_ALIGN.CENTER
    
    # 表格内容
    for row_idx, (commit_hash, date, commit_type, desc) in enumerate(commits):
        y_pos = 2.4 + row_idx * 0.35
        
        # 交替行背景
        if row_idx % 2 == 0:
            row_bg = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0.5), Inches(y_pos),
                Inches(9.0), Inches(0.35)
            )
            row_bg.fill.solid()
            row_bg.fill.fore_color.rgb = RGBColor(0xF8, 0xF9, 0xFA)
            row_bg.line.fill.background()
        
        # 单元格内容
        cells = [commit_hash, date, commit_type, desc]
        for col_idx, cell_text in enumerate(cells):
            cell_box = slide.shapes.add_textbox(Inches(col_x[col_idx]), Inches(y_pos), 
                                                 Inches(col_widths[col_idx]), Inches(0.35))
            tf = cell_box.text_frame
            tf.text = cell_text
            p = tf.paragraphs[0]
            p.font.size = Pt(8)
            p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
            
            # feat标绿，fix标红
            if col_idx == 2:
                if cell_text == "feat":
                    p.font.color.rgb = RGBColor(0x10, 0xB9, 0x81)
                elif cell_text == "fix":
                    p.font.color.rgb = RGBColor(0xEF, 0x44, 0x44)


def add_bug_fix_card(slide, x, y, width, height, bug_id, title, problem, solution, commit, color):
    """添加Bug修复卡片"""
    # 卡片背景
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y),
        Inches(width), Inches(height)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    card.line.color.rgb = color
    card.line.width = Pt(1.5)
    
    # Bug ID和标题
    title_box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.1), 
                                          Inches(3.0), Inches(0.3))
    tf = title_box.text_frame
    tf.text = bug_id + ": " + title
    p = tf.paragraphs[0]
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = color
    
    # 问题
    problem_box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.45), 
                                            Inches(width - 0.4), Inches(0.25))
    tf2 = problem_box.text_frame
    tf2.text = "问题: " + problem
    p2 = tf2.paragraphs[0]
    p2.font.size = Pt(9)
    p2.font.color.rgb = RGBColor(0xEF, 0x44, 0x44)
    
    # 解决方案
    solution_box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.72), 
                                            Inches(width - 0.4), Inches(0.25))
    tf3 = solution_box.text_frame
    tf3.text = "解决: " + solution
    p3 = tf3.paragraphs[0]
    p3.font.size = Pt(9)
    p3.font.color.rgb = RGBColor(0x10, 0xB9, 0x81)
    
    # 提交Hash
    commit_box = slide.shapes.add_textbox(Inches(x + width - 1.5), Inches(y + 0.1), 
                                           Inches(1.3), Inches(0.3))
    tf4 = commit_box.text_frame
    tf4.text = "commit: " + commit
    p4 = tf4.paragraphs[0]
    p4.font.size = Pt(8)
    p4.font.color.rgb = RGBColor(0x6B, 0x72, 0x80)
    p4.alignment = PP_ALIGN.RIGHT


def add_deployment_architecture(slide, color1, color2):
    """添加部署架构图"""
    services = [
        ("frontend", "Frontend\n(Nginx Alpine)", color1, 0.5, 1.5),
        ("backend", "Backend\n(Python 3.12)", color1, 0.5, 2.5),
        ("postgres", "PostgreSQL\n(16-alpine)", color2, 3.5, 1.5),
        ("redis", "Redis\n(7-alpine)", color2, 3.5, 2.5),
        ("chromadb", "ChromaDB\n(latest)", color2, 6.5, 1.5),
        ("arq-worker", "arq-worker\n(异步任务)", color1, 6.5, 2.5)
    ]
    
    for service_id, service_name, color, x, y in services:
        # 服务框
        service_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y),
            Inches(2.5), Inches(0.8)
        )
        service_box.fill.solid()
        service_box.fill.fore_color.rgb = color
        service_box.line.fill.background()
        
        tf = service_box.text_frame
        tf.text = service_name
        p = tf.paragraphs[0]
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.alignment = PP_ALIGN.CENTER


if __name__ == "__main__":
    try:
        output_file = create_professional_presentation()
        print("\n🎉 专业版PPT文件已生成：{}".format(output_file))
    except Exception as e:
        print("❌ 生成PPT时出错：{}".format(e))
        import traceback
        traceback.print_exc()
