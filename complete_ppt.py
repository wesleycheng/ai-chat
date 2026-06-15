#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
AI Chat Platform - 完整版项目文档 PPT 生成器
包含：项目概述、系统架构、技术栈、功能模块、开发过程、
      提交日志、Bug修复、部署架构、API接口文档、代码仓库信息
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE


def create_complete_presentation():
    """创建完整版项目演示 PPT"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # 定义专业配色方案
    COLOR_PRIMARY = RGBColor(0x1F, 0x4F, 0x8F)      # 深蓝色
    COLOR_SECONDARY = RGBColor(0x3B, 0x82, 0xF6)    # 亮蓝色
    COLOR_ACCENT = RGBColor(0xF5, 0x93, 0x1E)       # 橙色
    COLOR_SUCCESS = RGBColor(0x10, 0xB9, 0x81)      # 绿色
    COLOR_DARK = RGBColor(0x1F, 0x29, 0x3B)        # 深灰色
    COLOR_LIGHT_GRAY = RGBColor(0xF1, 0xF5, 0xF9)  # 浅灰色
    COLOR_TEXT = RGBColor(0x33, 0x33, 0x33)         # 正文灰色
    COLOR_PURPLE = RGBColor(0x8B, 0x5C, 0xF6)      # 紫色
    
    # ========== 幻灯片 1: 专业封面 ==========
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    add_professional_title_slide(slide1, "AI Chat Platform",
                                  "Enterprise Multi-Modal AI Chat Platform\n项目开发与技术文档",
                                  COLOR_PRIMARY, COLOR_SECONDARY)
    add_text_box(slide1, "Prepared by: Wesley Cheng", 7.5, 6.8, 2.0, 0.4,
                  size=14, color=COLOR_DARK, align=PP_ALIGN.RIGHT)
    add_text_box(slide1, "Date: May 30, 2026", 7.5, 7.1, 2.0, 0.4,
                  size=14, color=COLOR_DARK, align=PP_ALIGN.RIGHT)
    add_text_box(slide1, "Version: v2.0 (Complete Edition)", 0.5, 6.8, 3.0, 0.4,
                  size=12, color=COLOR_ACCENT, align=PP_ALIGN.LEFT)
    
    # ========== 幻灯片 2: 项目概述 ==========
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide2, "项目概述", COLOR_PRIMARY, COLOR_LIGHT_GRAY)
    add_info_card(slide2, 0.5, 1.2, 4.5, 2.8, "📊 项目信息",
                  ["项目名称: AI Chat Platform",
                   "开发周期: 2026-05-27 ~ 2026-05-30",
                   "总提交次数: 28 commits",
                   "代码规模: 前端~15K行 / 后端~8K行",
                   "生产环境: http://8.137.103.202",
                   "GitHub: github.com/wesleycheng/ai-chat"],
                  COLOR_SECONDARY)
    add_info_card(slide2, 5.2, 1.2, 4.3, 2.8, "✨ 核心特性",
                  ["✓ 多模型支持 (DeepSeek/OpenAI/Anthropic/Ollama)",
                   "✓ 智能文件解析 (PDF/DOCX/XLSX/TXT/MD)",
                   "✓ Agent管理系统 (自定义Prompt与工具)",
                   "✓ 流式对话 (SSE实时输出)",
                   "✓ 响应式设计 (桌面端+移动端)",
                   "✓ 安全认证 (JWT + bcrypt)"],
                  COLOR_SUCCESS)
    add_highlight_bar(slide2, 0.5, 4.5, 9.0, 0.8,
                      "技术亮点: 前后端分离 | Docker容器化 | RAG检索增强 | 异步任务队列",
                      COLOR_ACCENT)
    add_text_box(slide2, "GitHub仓库: https://github.com/wesleycheng/ai-chat", 
                  0.5, 5.5, 9.0, 0.4, size=11, color=COLOR_PURPLE, align=PP_ALIGN.CENTER)
    
    # ========== 幻灯片 3: 系统架构 ==========
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide3, "系统架构设计", COLOR_PRIMARY, COLOR_LIGHT_GRAY)
    add_architecture_diagram(slide3, COLOR_PRIMARY, COLOR_SECONDARY, COLOR_SUCCESS)
    add_text_box(slide3, "架构特点: 前后端分离 | RESTful API | SSE流式传输 | 容器化部署", 
                  0.5, 5.8, 9.0, 0.5, size=10, color=COLOR_DARK, align=PP_ALIGN.CENTER)
    
    # ========== 幻灯片 4: 技术栈 ==========
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide4, "技术栈", COLOR_PRIMARY, COLOR_LIGHT_GRAY)
    add_tech_stack_card(slide4, 0.5, 1.2, 4.5, 5.0, "前端技术栈",
                        ["React 18.3 + TypeScript 5.5",
                         "Vite 5.4 (构建工具)",
                         "Tailwind CSS 3.4 (样式)",
                         "Zustand 4.5 (状态管理)",
                         "@tanstack/react-query 5.56",
                         "React Router DOM 6.26",
                         "Axios 1.7 (HTTP客户端)",
                         "react-markdown 9.0",
                         "lucide-react 0.400 (图标)"],
                        COLOR_SECONDARY)
    add_tech_stack_card(slide4, 5.2, 1.2, 4.3, 5.0, "后端技术栈",
                        ["FastAPI 0.115 + Uvicorn",
                         "SQLAlchemy 2.0 + Alembic",
                         "asyncpg 0.30 + psycopg2-binary",
                         "python-jose 3.3 + passlib 1.7",
                         "LangChain 0.3 + LangGraph 0.2",
                         "ChromaDB 0.5 (向量数据库)",
                         "arq 0.24 + Redis 5.2",
                         "PyMuPDF 1.25 + python-docx 1.1"],
                        COLOR_SUCCESS)
    
    # ========== 幻灯片 5: 功能模块 ==========
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide5, "功能模块", COLOR_PRIMARY, COLOR_LIGHT_GRAY)
    modules = [
        ("1", "用户认证模块", ["JWT Token认证", "用户注册/登录", "密码bcrypt加密", "Token刷新机制"]),
        ("2", "模型配置模块", ["多模型支持", "CRUD管理", "API Key加密", "默认模型设置"]),
        ("3", "对话管理模块", ["SSE流式输出", "上下文管理", "Agent选择", "多轮对话"]),
        ("4", "文件管理模块", ["多格式支持", "异步上传", "自动解析", "RAG集成"]),
        ("5", "Agent管理模块", ["自定义Prompt", "工具配置", "CRUD管理", "模型关联"])
    ]
    x_positions = [0.5, 2.3, 4.1, 5.9, 7.7]
    for i, (num, title, features) in enumerate(modules):
        add_module_card(slide5, x_positions[i], 1.5, 1.6, 4.5,
                        num, title, features, COLOR_SECONDARY if i % 2 == 0 else COLOR_SUCCESS)
    
    # ========== 幻灯片 6: 开发过程 ==========
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide6, "开发过程", COLOR_PRIMARY, COLOR_LIGHT_GRAY)
    add_timeline(slide6,
                 [(0.8, "Day 1\n05-27", ["项目初始化", "Git仓库搭建", "Docker环境配置", "前后端骨架"]),
                  (2.8, "Day 2\n05-28", ["用户认证系统", "模型配置CRUD", "SSE流式输出", "ECS部署"]),
                  (4.8, "Day 3\n05-29", ["Agent管理页面", "文件上传功能", "Agent对话集成", "响应式优化"]),
                  (6.8, "Day 4\n05-30", ["默认模型编辑", "输入优化", "加载动画", "文档生成"])],
                 COLOR_SECONDARY)
    
    # ========== 幻灯片 7: 提交日志 ==========
    slide7 = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide7, "提交日志统计", COLOR_PRIMARY, COLOR_LIGHT_GRAY)
    add_commit_statistics(slide7, COLOR_SUCCESS)
    commits = [
        ("907c52b", "2026-05-30", "feat", "设置页面支持编辑默认模型和启用状态"),
        ("b23c987", "2026-05-29", "fix", "修复文件上传问题，修复认证头拼写错误"),
        ("eedd3dd", "2026-05-29", "feat", "优化所有页面响应式设计"),
        ("adcb7fb", "2026-05-29", "feat", "响应式UI优化"),
        ("c6a6c10", "2026-05-29", "fix", "修复文件上传问题"),
        ("19dba34", "2026-05-29", "fix", "修复文件内容获取问题"),
        ("fbf714c", "2026-05-29", "fix", "优化文件上传流程"),
        ("29ec3ab", "2026-05-29", "feat", "聊天支持上传文件"),
        ("0b76cfd", "2026-05-29", "feat", "聊天页面支持选择Agent"),
        ("9e72e15", "2026-05-29", "fix", "修复创建Agent时空字符串错误")
    ]
    add_commit_table(slide7, commits, COLOR_SECONDARY)
    
    # ========== 幻灯片 8: Bug修复记录 ==========
    slide8 = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide8, "重大Bug修复记录", COLOR_PRIMARY, COLOR_LIGHT_GRAY)
    bugs = [
        ("Bug #1", "Authorization拼写错误",
         "前端api.ts中Authorization拼写为Authoriation",
         "修正拼写错误", "b23c987", COLOR_ACCENT),
        ("Bug #2", "文件上传失败",
         "手动设置Content-Type导致boundary缺失",
         "删除手动Content-Type设置", "c6a6c10", COLOR_ACCENT),
        ("Bug #3", "Agent创建500错误",
         "空字符串''导致数据库外键约束报错",
         "将空字符串转为None", "9e72e15", COLOR_ACCENT),
        ("Bug #4", "文件内容未随消息发送",
         "file_ids类型不匹配，查询条件错误",
         "转换file_ids为UUID列表", "19dba34", COLOR_ACCENT)
    ]
    for i, (bug_id, title, problem, solution, commit, color) in enumerate(bugs):
        add_bug_fix_card(slide8, 0.5, 1.3 + i*1.3, 9.0, 1.1,
                         bug_id, title, problem, solution, commit, color)
    
    # ========== 幻灯片 9: 部署架构 ==========
    slide9 = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide9, "部署架构", COLOR_PRIMARY, COLOR_LIGHT_GRAY)
    add_deployment_architecture(slide9, COLOR_SECONDARY, COLOR_SUCCESS)
    add_info_card(slide9, 0.5, 5.5, 9.0, 1.5, "环境变量配置",
                  ["DATABASE_URL - 数据库连接字符串",
                   "REDIS_URL - Redis连接字符串",
                   "CHROMA_HOST - ChromaDB主机地址",
                   "SECRET_KEY - JWT密钥",
                   "GitHub仓库: https://github.com/wesleycheng/ai-chat"],
                  COLOR_ACCENT)
    
    # ========== 幻灯片 10: API接口概览 ==========
    slide10 = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide10, "API接口概览", COLOR_PRIMARY, COLOR_LIGHT_GRAY)
    add_info_card(slide10, 0.5, 1.2, 9.0, 1.0, "📡 接口统计",
                  ["总接口数: 27个 | 需要认证: 25个 | 无需认证: 2个（注册、登录） | 基础URL: http://8.137.103.202/api"],
                  COLOR_SECONDARY)
    
    api_modules = [
        ("认证模块 (4个)", ["POST /auth/register - 注册", 
                           "POST /auth/login - 登录",
                           "POST /auth/refresh - 刷新Token",
                           "GET /auth/me - 获取当前用户"]),
        ("对话模块 (6个)", ["GET/POST /conversations - 列表/创建",
                           "GET/DELETE /conversations/{id} - 详情/删除",
                           "GET /conversations/{id}/messages - 消息历史",
                           "POST /conversations/{id}/chat - 发送消息"]),
        ("文件模块 (4个)", ["POST /files/upload - 上传文件",
                           "GET /files - 文件列表",
                           "GET /files/{id} - 文件详情",
                           "DELETE /files/{id} - 删除文件"]),
        ("配置模块 (5个)", ["GET/POST /config/models - 列表/创建",
                           "PUT/DELETE /config/models/{id} - 更新/删除",
                           "POST /config/models/{id}/test - 测试模型"]),
        ("Agent模块 (6个)", ["GET/POST /agents - 列表/创建",
                            "GET/PUT/DELETE /agents/{id} - 详情/更新/删除",
                            "POST /agents/{id}/test - 测试Agent"])
    ]
    
    x_pos = [0.5, 2.3, 4.1, 5.9, 7.7]
    for i, (module_name, apis) in enumerate(api_modules):
        add_api_module_card(slide10, x_pos[i], 2.5, 1.6, 4.0,
                           module_name, apis, COLOR_SUCCESS if i % 2 == 0 else COLOR_PURPLE)
    
    # ========== 幻灯片 11: 核心API详解 ==========
    slide11 = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide11, "核心API详解", COLOR_PRIMARY, COLOR_LIGHT_GRAY)
    
    # 聊天接口详解
    add_info_card(slide11, 0.5, 1.2, 9.0, 1.5, "📨 POST /conversations/{id}/chat - 发送消息",
                  ["功能: 发送消息并获取AI回复（支持SSE流式输出）",
                   "请求体: {content: string, model_id: string, file_ids: string[], stream: boolean}",
                   "响应: SSE流 (data: {content: chunk}) 或 JSON {content: full_reply}",
                   "说明: 自动保存用户消息和AI回复，支持文件结合对话"],
                  COLOR_SECONDARY)
    
    # 文件上传接口详解
    add_info_card(slide11, 0.5, 3.0, 9.0, 1.5, "📎 POST /files/upload - 上传文件",
                  ["功能: 上传文件并立即解析内容（支持PDF/DOCX/XLSX/TXT/MD/JPG/PNG）",
                   "请求: multipart/form-data (file字段)",
                   "限制: 文件大小≤10MB",
                   "响应: {id, filename, file_ext, content_text, parse_status}"],
                  COLOR_SUCCESS)
    
    # 模型配置接口详解
    add_info_card(slide11, 0.5, 4.8, 9.0, 1.5, "⚙️ POST /config/models - 创建模型配置",
                  ["功能: 创建新的AI模型配置（API Key会加密存储）",
                   "请求体: {name, provider, api_base, api_key, model_name, is_default}",
                   "说明: is_default=true时会自动取消其他模型的默认状态",
                   "安全: API Key使用AES-256加密，返回时自动脱敏"],
                  COLOR_ACCENT)
    
    # ========== 幻灯片 12: 代码仓库信息 ==========
    slide12 = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide12, "代码仓库信息", COLOR_PRIMARY, COLOR_LIGHT_GRAY)
    
    add_info_card(slide12, 0.5, 1.2, 9.0, 1.8, "📦 GitHub仓库",
                  ["仓库地址: <ADDRESS_REMOVED> <EMAIL_REMOVED>",
                   "克隆方式: git clone git@github.com:wesleycheng/ai-chat.git",
                   "生产环境: http://8.137.103.202",
                   "在线API文档: http://8.137.103.202/api/docs"],
                  COLOR_SECONDARY)
    
    add_info_card(slide12, 0.5, 3.2, 4.5, 3.5, "📂 项目结构",
                  ["ai-chat/",
                   "├── frontend/  (React + TypeScript)",
                   "├── backend/   (FastAPI)",
                   "├── nginx/     (Nginx配置)",
                   "├── docker-compose.yml",
                   "├── deploy-ecs.sh",
                   "└── PROJECT_DOCUMENTATION.md"],
                  COLOR_SUCCESS)
    
    add_info_card(slide12, 5.2, 3.2, 4.3, 3.5, "🔧 开发环境",
                  ["前置要求:",
                   "  • Python 3.12+",
                   "  • Node.js 20+",
                   "  • Docker Desktop",
                   "  • Git",
                   "",
                   "快速启动: docker compose up -d"],
                  COLOR_PURPLE)
    
    # ========== 幻灯片 13: 部署指南 ==========
    slide13 = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide13, "部署指南", COLOR_PRIMARY, COLOR_LIGHT_GRAY)
    
    add_info_card(slide13, 0.5, 1.2, 9.0, 2.0, "🚀 ECS生产环境部署",
                  ["服务器: 8.137.103.202 (Ubuntu 22.04 LTS)",
                   "部署步骤:",
                   "  1. SSH登录: ssh root@8.137.103.202",
                   "  2. 进入目录: cd /var/www/ai-chat",
                   "  3. 拉取代码: git pull origin main",
                   "  4. 执行部署: bash deploy-ecs.sh",
                   "  5. 检查状态: docker compose ps"],
                  COLOR_SECONDARY)
    
    add_info_card(slide13, 0.5, 3.5, 4.5, 3.0, "🐳 Docker服务",
                  ["frontend  - Nginx Alpine (前端)",
                   "backend   - Python 3.12 (后端)",
                   "postgres  - PostgreSQL 16 (数据库)",
                   "redis     - Redis 7 (缓存)",
                   "chromadb  - ChromaDB (向量数据库)",
                   "arq-worker - 异步任务队列"],
                  COLOR_SUCCESS)
    
    add_info_card(slide13, 5.2, 3.5, 4.3, 3.0, "🔐 环境变量",
                  ["DATABASE_URL  - 数据库连接",
                   "REDIS_URL     - Redis连接",
                   "SECRET_KEY    - JWT密钥",
                   "CHROMA_HOST   - ChromaDB地址",
                   "UPLOAD_DIR    - 文件上传目录",
                   "CORS_ORIGINS  - 跨域配置"],
                  COLOR_ACCENT)
    
    # ========== 幻灯片 14: 使用指南 ==========
    slide14 = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide14, "使用指南 & 未来规划", COLOR_PRIMARY, COLOR_LIGHT_GRAY)
    add_info_card(slide14, 0.5, 1.2, 4.5, 2.8, "📖 快速开始",
                  ["1. 访问 http://8.137.103.202",
                   "2. 注册账号并登录",
                   "3. 配置模型（设置->添加模型）",
                   "4. 开始对话（选择模型/Agent）",
                   "5. 上传文件进行RAG对话",
                   "",
                   "详细文档: 见PROJECT_DOCUMENTATION.md"],
                  COLOR_SECONDARY)
    add_info_card(slide14, 5.2, 1.2, 4.3, 2.8, "🚀 未来规划",
                  ["功能增强:",
                   "  • 支持更多文件类型(PPT/MP4)",
                   "  • 实现RAG检索增强",
                   "  • 多模态输入(图片/语音)",
                   "性能优化:",
                   "  • 消息分页加载",
                   "  • 文件解析速度优化",
                   "安全加固:",
                   "  • API限流",
                   "  • 审计日志"],
                  COLOR_SUCCESS)
    add_highlight_bar(slide14, 0.5, 4.5, 9.0, 0.8,
                      "项目已部署生产环境，欢迎测试与使用！GitHub: https://github.com/wesleycheng/ai-chat",
                      COLOR_ACCENT)
    
    # ========== 幻灯片 15: 致谢 ==========
    slide15 = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 背景
    bg = slide15.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                                   Inches(10), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLOR_PRIMARY
    bg.line.fill.background()
    
    # 标题
    title_box = slide15.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.0))
    tf = title_box.text_frame
    tf.text = "Thank You!"
    p = tf.paragraphs[0]
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.CENTER
    
    # 副标题
    subtitle_box = slide15.shapes.add_textbox(Inches(1), Inches(3.8), Inches(8), Inches(0.8))
    tf2 = subtitle_box.text_frame
    tf2.text = "AI Chat Platform - Enterprise Multi-Modal AI Chat Platform"
    p2 = tf2.paragraphs[0]
    p2.font.size = Pt(16)
    p2.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    p2.alignment = PP_ALIGN.CENTER
    
    # 联系信息
    contact_box = slide15.shapes.add_textbox(Inches(1), Inches(5.0), Inches(8), Inches(1.0))
    tf3 = contact_box.text_frame
    tf3.text = "Developer: Wesley Cheng\nGitHub: https://github.com/wesleycheng\nEmail: [Your Email]"
    p3 = tf3.paragraphs[0]
    p3.font.size = Pt(12)
    p3.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    p3.alignment = PP_ALIGN.CENTER
    
    # 保存文件
    output_path = "/Users/wesleycheng/Documents/ai-chat-platform/AI_Chat_Platform_完整版项目文档.pptx"
    prs.save(output_path)
    print("✅ 完整版PPT生成成功！文件路径：{}".format(output_path))
    print("📊 PPT包含15页幻灯片，涵盖所有项目内容")
    return output_path


# ============ 辅助函数 ============

def add_professional_title_slide(slide, title, subtitle, color_primary, color_secondary):
    """添加专业封面"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                                    Inches(10), Inches(7.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xF8, 0xFA, 0xFC)
    shape.line.fill.background()
    
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                                      Inches(10), Inches(0.1))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = color_primary
    top_bar.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.0), Inches(8), Inches(1.0))
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = color_primary
    p.alignment = PP_ALIGN.CENTER
    
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.2), Inches(8), Inches(1.2))
    tf2 = subtitle_box.text_frame
    tf2.text = subtitle
    p2 = tf2.paragraphs[0]
    p2.font.size = Pt(18)
    p2.font.color.rgb = RGBColor(0x64, 0x74, 0x8B)
    p2.alignment = PP_ALIGN.CENTER


def add_slide_header(slide, title, color_primary, color_bg):
    """添加幻灯片标题栏"""
    header_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                                        Inches(10), Inches(0.8))
    header_bar.fill.solid()
    header_bar.fill.fore_color.rgb = color_primary
    header_bar.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.LEFT


def add_info_card(slide, x, y, width, height, title, items, color):
    """添加信息卡片"""
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y),
                                  Inches(width), Inches(height))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    card.line.color.rgb = color
    card.line.width = Pt(1.5)
    
    title_box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.15),
                                          Inches(width - 0.4), Inches(0.4))
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = color
    
    content_y = y + 0.6
    for item in items:
        item_box = slide.shapes.add_textbox(Inches(x + 0.3), Inches(content_y),
                                             Inches(width - 0.6), Inches(0.35))
        tf2 = item_box.text_frame
        tf2.text = item
        p2 = tf2.paragraphs[0]
        p2.font.size = Pt(10)
        p2.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        content_y += 0.35


def add_text_box(slide, text, x, y, width, height, size=12,
                 color=RGBColor(0, 0, 0), align=PP_ALIGN.LEFT):
    """添加文本框"""
    text_box = slide.shapes.add_textbox(Inches(x), Inches(y),
                                         Inches(width), Inches(height))
    tf = text_box.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.alignment = align


def add_highlight_bar(slide, x, y, width, height, text, color):
    """添加高亮条"""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
                                  Inches(width), Inches(height))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    
    text_box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.1),
                                         Inches(width - 0.4), Inches(height - 0.2))
    tf = text_box.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.CENTER


def add_architecture_diagram(slide, color1, color2, color3):
    """添加架构图"""
    add_arch_box(slide, 3.5, 1.2, 3.0, 0.6, "用户层 (Web Browser)", color1)
    add_connector(slide, 5.0, 1.8, 5.0, 2.2, color1)
    add_arch_box(slide, 2.5, 2.2, 5.0, 0.6, "前端层 (React + TypeScript)", color2)
    add_connector(slide, 5.0, 2.8, 5.0, 3.2, color2)
    add_arch_box(slide, 2.5, 3.2, 5.0, 0.6, "后端层 (FastAPI + Uvicorn)", color2)
    add_connector(slide, 5.0, 3.8, 5.0, 4.2, color2)
    add_arch_box(slide, 2.5, 4.2, 5.0, 0.6, "数据层 (PostgreSQL + Redis + ChromaDB)", color3)
    
    features = ["前后端分离", "Docker容器化", "异步任务(arq)", "向量检索(ChromaDB)", "流式输出(SSE)"]
    for i, feature in enumerate(features):
        feature_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                             Inches(7.8), Inches(1.5 + i * 0.7),
                                             Inches(1.8), Inches(0.5))
        feature_box.fill.solid()
        feature_box.fill.fore_color.rgb = color3
        feature_box.line.fill.background()
        
        tf = feature_box.text_frame
        tf.text = feature
        p = tf.paragraphs[0]
        p.font.size = Pt(9)
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.alignment = PP_ALIGN.CENTER


def add_arch_box(slide, x, y, width, height, text, color):
    """添加架构框"""
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y),
                                  Inches(width), Inches(height))
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.line.color.rgb = color
    box.line.width = Pt(2)
    
    tf = box.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.CENTER


def add_connector(slide, x1, y1, x2, y2, color):
    """添加连接线"""
    connector = slide.shapes.add_connector(1, Inches(x1), Inches(y1),
                                           Inches(x2), Inches(y2))
    connector.line.color.rgb = color
    connector.line.width = Pt(2)


def add_tech_stack_card(slide, x, y, width, height, title, tech_list, color):
    """添加技术栈卡片"""
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y),
                                  Inches(width), Inches(height))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    card.line.color.rgb = color
    card.line.width = Pt(2)
    
    title_box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.15),
                                          Inches(width - 0.4), Inches(0.4))
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = color
    
    content_y = y + 0.6
    for tech in tech_list:
        tech_box = slide.shapes.add_textbox(Inches(x + 0.25), Inches(content_y),
                                             Inches(width - 0.5), Inches(0.32))
        tf2 = tech_box.text_frame
        tf2.text = "• " + tech
        p2 = tf2.paragraphs[0]
        p2.font.size = Pt(9)
        p2.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        content_y += 0.32


def add_module_card(slide, x, y, width, height, num, title, features, color):
    """添加功能模块卡片"""
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y),
                                  Inches(width), Inches(height))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    card.line.color.rgb = color
    card.line.width = Pt(1.5)
    
    num_circle = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                         Inches(x + width/2 - 0.25), Inches(y + 0.15),
                                         Inches(0.5), Inches(0.5))
    num_circle.fill.solid()
    num_circle.fill.fore_color.rgb = color
    num_circle.line.fill.background()
    
    tf_num = num_circle.text_frame
    tf_num.text = num
    p_num = tf_num.paragraphs[0]
    p_num.font.size = Pt(16)
    p_num.font.bold = True
    p_num.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p_num.alignment = PP_ALIGN.CENTER
    
    title_box = slide.shapes.add_textbox(Inches(x + 0.1), Inches(y + 0.75),
                                          Inches(width - 0.2), Inches(0.35))
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.CENTER
    
    content_y = y + 1.15
    for feature in features:
        feat_box = slide.shapes.add_textbox(Inches(x + 0.15), Inches(content_y),
                                             Inches(width - 0.3), Inches(0.28))
        tf2 = feat_box.text_frame
        tf2.text = "✓ " + feature
        p2 = tf2.paragraphs[0]
        p2.font.size = Pt(8)
        p2.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        content_y += 0.28


def add_api_module_card(slide, x, y, width, height, module_name, apis, color):
    """添加API模块卡片"""
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y),
                                  Inches(width), Inches(height))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    card.line.color.rgb = color
    card.line.width = Pt(1.5)
    
    title_box = slide.shapes.add_textbox(Inches(x + 0.1), Inches(y + 0.1),
                                          Inches(width - 0.2), Inches(0.3))
    tf = title_box.text_frame
    tf.text = module_name
    p = tf.paragraphs[0]
    p.font.size = Pt(9)
    p.font.bold = True
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.CENTER
    
    content_y = y + 0.45
    for api in apis:
        api_box = slide.shapes.add_textbox(Inches(x + 0.15), Inches(content_y),
                                            Inches(width - 0.3), Inches(0.55))
        tf2 = api_box.text_frame
        tf2.text = api
        p2 = tf2.paragraphs[0]
        p2.font.size = Pt(7)
        p2.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        content_y += 0.55


def add_timeline(slide, milestones, color):
    """添加时间线"""
    for i, (x, label, items) in enumerate(milestones):
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(1.5),
                                         Inches(0.6), Inches(0.6))
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.line.fill.background()
        
        label_box = slide.shapes.add_textbox(Inches(x - 0.2), Inches(1.5),
                                              Inches(1.0), Inches(0.6))
        tf = label_box.text_frame
        tf.text = label
        p = tf.paragraphs[0]
        p.font.size = Pt(9)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.alignment = PP_ALIGN.CENTER
        
        if i < len(milestones) - 1:
            connector = slide.shapes.add_connector(1, Inches(x + 0.6), Inches(1.8),
                                                   Inches(milestones[i+1][0]), Inches(1.8))
            connector.line.color.rgb = color
            connector.line.width = Pt(2)
        
        content_y = 2.3
        for item in items:
            item_box = slide.shapes.add_textbox(Inches(x - 0.3), Inches(content_y),
                                                 Inches(1.8), Inches(0.35))
            tf2 = item_box.text_frame
            tf2.text = "✅ " + item
            p2 = tf2.paragraphs[0]
            p2.font.size = Pt(8)
            p2.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
            content_y += 0.35


def add_commit_statistics(slide, color):
    """添加提交统计"""
    stats = [("✨ feat", "15次", RGBColor(0x10, 0xB9, 0x81)),
             ("🐛 fix", "12次", RGBColor(0xEF, 0x44, 0x44)),
             ("📝 docs", "1次", RGBColor(0x8B, 0x5C, 0xF6)),
             ("🔧 chore", "0次", RGBColor(0xF5, 0x93, 0x1E))]
    
    for i, (label, count, color_stat) in enumerate(stats):
        stat_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                           Inches(0.5 + i * 2.3), Inches(1.2),
                                           Inches(2.0), Inches(0.6))
        stat_box.fill.solid()
        stat_box.fill.fore_color.rgb = color_stat
        stat_box.line.fill.background()
        
        tf = stat_box.text_frame
        tf.text = label + ": " + count
        p = tf.paragraphs[0]
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.alignment = PP_ALIGN.CENTER


def add_commit_table(slide, commits, color):
    """添加提交记录表格"""
    headers = ["Commit Hash", "日期", "类型", "描述"]
    col_widths = [1.5, 1.2, 0.8, 4.5]
    col_x = [0.5, 2.0, 2.8, 3.6]
    
    header_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(2.0),
                                        Inches(9.0), Inches(0.4))
    header_bg.fill.solid()
    header_bg.fill.fore_color.rgb = color
    header_bg.line.fill.background()
    
    for i, header in enumerate(headers):
        header_box = slide.shapes.add_textbox(Inches(col_x[i]), Inches(2.0),
                                              Inches(col_widths[i]), Inches(0.4))
        tf = header_box.text_frame
        tf.text = header
        p = tf.paragraphs[0]
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.alignment = PP_ALIGN.CENTER
    
    for row_idx, (commit_hash, date, commit_type, desc) in enumerate(commits):
        y_pos = 2.4 + row_idx * 0.35
        
        if row_idx % 2 == 0:
            row_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(y_pos),
                                            Inches(9.0), Inches(0.35))
            row_bg.fill.solid()
            row_bg.fill.fore_color.rgb = RGBColor(0xF8, 0xF9, 0xFA)
            row_bg.line.fill.background()
        
        cells = [commit_hash, date, commit_type, desc]
        for col_idx, cell_text in enumerate(cells):
            cell_box = slide.shapes.add_textbox(Inches(col_x[col_idx]), Inches(y_pos),
                                                 Inches(col_widths[col_idx]), Inches(0.35))
            tf = cell_box.text_frame
            tf.text = cell_text
            p = tf.paragraphs[0]
            p.font.size = Pt(8)
            p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
            
            if col_idx == 2:
                if cell_text == "feat":
                    p.font.color.rgb = RGBColor(0x10, 0xB9, 0x81)
                elif cell_text == "fix":
                    p.font.color.rgb = RGBColor(0xEF, 0x44, 0x44)


def add_bug_fix_card(slide, x, y, width, height, bug_id, title, problem, solution, commit, color):
    """添加Bug修复卡片"""
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y),
                                  Inches(width), Inches(height))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    card.line.color.rgb = color
    card.line.width = Pt(1.5)
    
    title_box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.1),
                                          Inches(3.0), Inches(0.3))
    tf = title_box.text_frame
    tf.text = bug_id + ": " + title
    p = tf.paragraphs[0]
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = color
    
    problem_box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.45),
                                            Inches(width - 0.4), Inches(0.25))
    tf2 = problem_box.text_frame
    tf2.text = "问题: " + problem
    p2 = tf2.paragraphs[0]
    p2.font.size = Pt(9)
    p2.font.color.rgb = RGBColor(0xEF, 0x44, 0x44)
    
    solution_box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.72),
                                            Inches(width - 0.4), Inches(0.25))
    tf3 = solution_box.text_frame
    tf3.text = "解决: " + solution
    p3 = tf3.paragraphs[0]
    p3.font.size = Pt(9)
    p3.font.color.rgb = RGBColor(0x10, 0xB9, 0x81)
    
    commit_box = slide.shapes.add_textbox(Inches(x + width - 1.5), Inches(y + 0.1),
                                           Inches(1.3), Inches(0.3))
    tf4 = commit_box.text_frame
    tf4.text = "commit: " + commit
    p4 = tf4.paragraphs[0]
    p4.font.size = Pt(8)
    p4.font.color.rgb = RGBColor(0x6B, 0x72, 0x80)
    p4.alignment = PP_ALIGN.RIGHT


def add_deployment_architecture(slide, color1, color2):
    """添加部署架构图"""
    services = [("frontend", "Frontend\n(Nginx Alpine)", color1, 0.5, 1.5),
                ("backend", "Backend\n(Python 3.12)", color1, 0.5, 2.5),
                ("postgres", "PostgreSQL\n(16-alpine)", color2, 3.5, 1.5),
                ("redis", "Redis\n(7-alpine)", color2, 3.5, 2.5),
                ("chromadb", "ChromaDB\n(latest)", color2, 6.5, 1.5),
                ("arq-worker", "arq-worker\n(异步任务)", color1, 6.5, 2.5)]
    
    for service_id, service_name, color, x, y in services:
        service_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                             Inches(x), Inches(y),
                                             Inches(2.5), Inches(0.8))
        service_box.fill.solid()
        service_box.fill.fore_color.rgb = color
        service_box.line.fill.background()
        
        tf = service_box.text_frame
        tf.text = service_name
        p = tf.paragraphs[0]
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.alignment = PP_ALIGN.CENTER


if __name__ == "__main__":
    try:
        output_file = create_complete_presentation()
        print("\n🎉 完整版PPT文件已生成：{}".format(output_file))
    except Exception as e:
        print("❌ 生成PPT时出错：{}".format(e))
        import traceback
        traceback.print_exc()
