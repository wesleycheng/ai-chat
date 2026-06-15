#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
AI Chat Platform - 项目文档 PPT 生成器
使用 python-pptx 生成项目演示文稿
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


def create_presentation():
    """创建项目演示 PPT"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # 幻灯片 1: 封面
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide1, "AI Chat Platform", "企业级多功能 AI 聊天平台\n\n项目开发文档", 1.5, 2.5)
    add_text(slide1, "开发者：Wesley Cheng\n日期：2026-05-30", 4.0, 5.5, size=14)
    
    # 幻灯片 2: 项目概述
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide2, "项目概述", "", 0.5, 0.5)
    content2 = (
        "✨ 核心特性\n\n"
        "• 多模型支持 - DeepSeek、OpenAI、Anthropic、Ollama、自定义\n"
        "• 文件智能解析 - PDF、DOCX、XLSX、TXT、MD、JPG、PNG\n"
        "• Agent 管理系统 - 自定义系统提示词、工具配置\n"
        "• 流式对话 - SSE 实时输出\n"
        "• 响应式设计 - 完美适配桌面端和移动端\n"
        "• 安全认证 - JWT + bcrypt 加密\n\n"
        "📊 项目信息\n\n"
        "• 开发时间：2026-05-27 至 2026-05-30\n"
        "• 总提交次数：28 次\n"
        "• 生产环境：http://8.137.103.202"
    )
    add_text(slide2, content2, 0.8, 1.5)
    
    # 幻灯片 3: 系统架构
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide3, "系统架构", "", 0.5, 0.5)
    content3 = (
        "整体架构：\n\n"
        "用户层 (Web Browser)\n"
        "    ↓\n"
        "前端层 (React + TypeScript)\n"
        "ChatPage | Settings | Files | Agents\n"
        "    ↓ HTTP/WebSocket\n"
        "后端层 (FastAPI + Uvicorn)\n"
        "Auth | Chat | Files | Agents | Config\n"
        "    ↓\n"
        "数据层 (PostgreSQL + Redis + ChromaDB)\n\n"
        "技术特点：\n"
        "• 前后端分离架构\n"
        "• 容器化部署 (Docker + Docker Compose)\n"
        "• 异步任务处理 (arq)\n"
        "• 向量检索 (ChromaDB + RAG)\n"
        "• 流式输出 (SSE)"
    )
    add_text(slide3, content3, 0.6, 1.5, size=11)
    
    # 幻灯片 4: 技术栈
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide4, "技术栈", "", 0.5, 0.5)
    content4 = (
        "前端技术栈：\n"
        "• 框架：React 18.3 + TypeScript 5.5\n"
        "• 构建工具：Vite 5.4\n"
        "• 样式方案：Tailwind CSS 3.4\n"
        "• 状态管理：Zustand 4.5\n"
        "• 数据请求：@tanstack/react-query 5.56\n"
        "• 路由管理：React Router DOM 6.26\n"
        "• HTTP 客户端：Axios 1.7\n"
        "• Markdown 渲染：react-markdown 9.0\n"
        "• 图标库：lucide-react 0.400\n\n"
        "后端技术栈：\n"
        "• Web 框架：FastAPI 0.115 + Uvicorn\n"
        "• ORM：SQLAlchemy 2.0 + Alembic\n"
        "• 数据库驱动：asyncpg 0.30 + psycopg2-binary\n"
        "• 认证：python-jose 3.3 + passlib 1.7\n"
        "• AI 框架：LangChain 0.3 + LangGraph 0.2\n"
        "• 向量数据库：ChromaDB 0.5\n"
        "• 异步任务：arq 0.24 + Redis 5.2"
    )
    add_text(slide4, content4, 0.6, 1.5, size=10)
    
    # 幻灯片 5: 功能模块
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide5, "功能模块", "", 0.5, 0.5)
    content5 = (
        "1. 用户认证模块\n"
        "   • 用户注册/登录 (JWT Token 认证)\n"
        "   • Token 刷新机制、密码 bcrypt 加密存储\n\n"
        "2. 模型配置模块\n"
        "   • 支持多种 AI 模型提供商\n"
        "   • 模型配置 CRUD（创建、编辑、删除）\n"
        "   • 设置默认模型、启用/禁用\n"
        "   • API Key 加密存储 (AES-256)\n\n"
        "3. 对话管理模块\n"
        "   • 创建新对话、切换模型\n"
        "   • 选择 Agent 进行对话\n"
        "   • SSE 流式输出\n"
        "   • 上下文管理 (Sliding Window)\n\n"
        "4. 文件管理模块\n"
        "   • 支持文件类型：PDF、DOCX、XLSX、TXT、MD、JPG、PNG\n"
        "   • 文件大小限制：10MB\n"
        "   • 异步文件上传、自动解析\n\n"
        "5. Agent 管理模块\n"
        "   • Agent 创建、编辑、删除\n"
        "   • 自定义系统提示词 (System Prompt)"
    )
    add_text(slide5, content5, 0.6, 1.5, size=11)
    
    # 幻灯片 6: 开发过程
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide6, "开发过程", "", 0.5, 0.5)
    content6 = (
        "开发时间线：\n\n"
        "Day 1 (2026-05-27) - 项目初始化\n"
        "  ✅ 初始化 Git 仓库、创建项目基础结构\n"
        "  ✅ 配置 Docker Compose 环境\n"
        "  ✅ 后端 FastAPI 骨架搭建\n"
        "  ✅ 前端 Vite + React + TypeScript 初始化\n\n"
        "Day 2 (2026-05-28) - 核心功能开发\n"
        "  ✅ 用户认证系统 (JWT + bcrypt)\n"
        "  ✅ 模型配置 CRUD、对话功能 + SSE 流式输出\n"
        "  ✅ 文件上传接口\n"
        "  ✅ 完善 Docker 部署配置\n"
        "  ✅ ECS 生产环境部署\n\n"
        "Day 3 (2026-05-29) - 功能完善\n"
        "  ✅ Agent 管理页面（创建、编辑、删除）\n"
        "  ✅ 聊天页面支持选择 Agent 进行对话\n"
        "  ✅ 聊天支持上传文件，结合文件内容进行对话\n"
        "  ✅ 全站响应式 UI 优化\n\n"
        "Day 4 (2026-05-30) - 细节优化\n"
        "  ✅ 设置页面支持编辑默认模型和启用状态\n"
        "  ✅ 优化输入内容过后显示问题\n"
        "  ✅ 等待答案过程中显示动画效果"
    )
    add_text(slide6, content6, 0.6, 1.5, size=10)
    
    # 幻灯片 7: 提交日志
    slide7 = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide7, "提交日志 (前10条)", "", 0.5, 0.5)
    content7 = (
        "907c52b | 2026-05-30 | feat: 设置页面支持编辑默认模型和启用状态\n"
        "b23c987 | 2026-05-29 | fix: 修复文件上传问题，修复认证头拼写错误\n"
        "eedd3dd | 2026-05-29 | feat: 优化所有页面响应式设计\n"
        "adcb7fb | 2026-05-29 | feat: 响应式UI优化\n"
        "c6a6c10 | 2026-05-29 | fix: 修复文件上传问题\n"
        "19dba34 | 2026-05-29 | fix: 修复文件内容获取问题\n"
        "fbf714c | 2026-05-29 | fix: 优化文件上传流程，修复上传失败问题\n"
        "29ec3ab | 2026-05-29 | feat: 聊天支持上传文件，结合文件内容进行对话\n"
        "0b76cfd | 2026-05-29 | feat: 聊天页面支持选择Agent进行对话\n"
        "9e72e15 | 2026-05-29 | fix: 修复创建Agent时model_id为空字符串导致500错误\n\n"
        "提交统计：\n"
        "• ✨ 新功能 (feat): 15 次\n"
        "• 🐛 错误修复 (fix): 12 次\n"
        "• 总计: 28 次提交"
    )
    add_text(slide7, content7, 0.6, 1.5, size=9)
    
    # 幻灯片 8: 修复与测试
    slide8 = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide8, "重大 Bug 修复", "", 0.5, 0.5)
    content8 = (
        "Bug #1: Authorization 拼写错误\n"
        "  • 问题：前端 api.ts 中 Authorization 拼写为 Authoriation\n"
        "  • 修复：修正拼写错误\n"
        "  • 提交: b23c987\n\n"
        "Bug #2: 文件上传失败（Content-Type 问题）\n"
        "  • 问题：手动设置 Content-Type 导致 boundary 缺失\n"
        "  • 修复：删除手动设置的 Content-Type\n"
        "  • 提交: c6a6c10\n\n"
        "Bug #3: 创建 Agent 时 model_id 为空字符串导致 500 错误\n"
        "  • 问题：空字符串 '' 导致数据库表键约束报错\n"
        "  • 修复：将空字符串转为 None\n"
        "  • 提交: 9e72e15\n\n"
        "Bug #4: 文件内容未随消息发送\n"
        "  • 问题：file_ids 类型不匹配，查询条件错误\n"
        "  • 修复：转换 file_ids 为 UUID 列表\n"
        "  • 提交: 19dba34\n\n"
        "测试情况：\n"
        "  ✅ 认证模块测试 - 通过\n"
        "  ✅ 模型配置测试 - 通过\n"
        "  ✅ 对话功能测试 - 通过\n"
        "  ✅ 文件上传测试 - 通过\n"
        "  ✅ Agent 管理测试 - 通过"
    )
    add_text(slide8, content8, 0.6, 1.5, size=10)
    
    # 幻灯片 9: 部署架构
    slide9 = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide9, "部署架构", "", 0.5, 0.5)
    content9 = (
        "Docker Compose 服务架构：\n\n"
        "• frontend (Nginx Alpine) - 端口：5173:80\n"
        "• backend (Python 3.12-slim) - 端口：8000:8000\n"
        "• postgres (PostgreSQL 16-alpine) - 端口：5432:5432\n"
        "• redis (Redis 7-alpine) - 端口：6379:6379\n"
        "• chromadb (ChromaDB latest) - 端口：8001:8000\n"
        "• arq-worker (异步任务队列)\n\n"
        "ECS 部署流程：\n"
        "  1. SSH 登录到 ECS 服务器\n"
        "  2. 进入项目目录：cd /var/www/ai-chat\n"
        "  3. 执行部署脚本：bash deploy-ecs.sh\n"
        "  4. 检查服务状态：docker compose ps\n\n"
        "环境变量配置：\n"
        "  • DATABASE_URL - 数据库连接字符串\n"
        "  • REDIS_URL - Redis 连接字符串\n"
        "  • CHROMA_HOST - ChromaDB 主机地址\n"
        "  • SECRET_KEY - JWT 密钥"
    )
    add_text(slide9, content9, 0.6, 1.5, size=10)
    
    # 幻灯片 10: 使用指南 & 未来规划
    slide10 = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide10, "使用指南 & 未来规划", "", 0.5, 0.5)
    content10 = (
        "快速开始：\n\n"
        "1. 注册账号\n"
        "   • 访问 http://8.137.103.202\n"
        "   • 点击"注册"按钮，输入用户名、密码\n\n"
        "2. 配置模型\n"
        "   • 登录后，点击"设置"\n"
        "   • 展开"模型配置"，点击"添加模型"\n"
        "   • 填写模型信息（名称、提供商、API Key）\n\n"
        "3. 开始对话\n"
        "   • 点击"新对话"\n"
        "   • 输入消息，按 Enter 发送\n\n"
        "未来规划：\n\n"
        "功能增强：\n"
        "  - 支持更多文件类型（PPT、MP4 等）\n"
        "  - 实现 RAG 检索增强（ChromaDB 向量存储）\n"
        "  - 支持多模态输入（图片、语音）\n\n"
        "性能优化：\n"
        "  - 实现消息分页加载\n"
        "  - 优化文件解析速度"
    )
    add_text(slide10, content10, 0.6, 1.5, size=10)
    
    # 保存文件
    output_path = "/Users/wesleycheng/Documents/ai-chat-platform/AI_Chat_Platform_项目文档.pptx"
    prs.save(output_path)
    print("✅ PPT 生成成功！文件路径：{}".format(output_path))
    return output_path


def add_title(slide, title_text, subtitle_text, x, y):
    """添加标题"""
    title_box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title_text
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0x1F, 0x29, 0x3B)
    title_para.alignment = PP_ALIGN.CENTER
    
    # 副标题（如果有）
    if subtitle_text:
        subtitle_box = slide.shapes.add_textbox(Inches(x), Inches(y + 0.8), Inches(9), Inches(0.5))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle_text
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.font.size = Pt(16)
        subtitle_para.font.color.rgb = RGBColor(0x6B, 0x72, 0x80)
        subtitle_para.alignment = PP_ALIGN.CENTER


def add_text(slide, text, x, y, size=12):
    """添加文本"""
    text_box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(9), Inches(5.5))
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    text_frame.text = text
    
    for paragraph in text_frame.paragraphs:
        paragraph.font.size = Pt(size)
        paragraph.font.color.rgb = RGBColor(0x37, 0x4A, 0x52)
        paragraph.space_after = Pt(6)


if __name__ == "__main__":
    try:
        output_file = create_presentation()
        print("\n🎉 PPT 文件已生成：{}".format(output_file))
    except ImportError as e:
        print("❌ 错误：未安装 python-pptx 库")
        print("请运行：pip install python-pptx")
        print("具体错误：{}".format(e))
    except Exception as e:
        print("❌ 生成 PPT 时出错：{}".format(e))
        import traceback
        traceback.print_exc()
